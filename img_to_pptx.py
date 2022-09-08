#モジュールのインポート
from pptx import Presentation
from pptx.util import Cm
import os
import glob
print(os.getcwd())
#プレゼンテーションオブジェクトの作成
ppt = Presentation()
w, h = 33.867, 19.05
ppt.slide_width = Cm(w)
ppt.slide_height = Cm(h)

problems = sorted(glob.glob(r'C:\Users\81901\working\kyopro\kyopro_educational_90\problem/*.jpg'))
sols = sorted(glob.glob(r'C:\Users\81901\working\kyopro\kyopro_educational_90\editorial/*.jpg'))

#タイトルスライドを追加
for p in problems:
    slide_layout_6 = ppt.slide_layouts[6]
    slide_6 = ppt.slides.add_slide(slide_layout_6)
    shapes = slide_6.shapes
    shapes.add_picture(p, 0, 0, width=Cm(w), height=Cm(h))
ppt.save(r'C:\Users\81901\working\kyopro\kyopro_educational_90/problems.pptx')

ppt = Presentation()
w, h = 33.867, 19.05
ppt.slide_width = Cm(w)
ppt.slide_height = Cm(h)

for s in sols:
    slide_layout_6 = ppt.slide_layouts[6]
    slide_6 = ppt.slides.add_slide(slide_layout_6)
    shapes = slide_6.shapes
    shapes.add_picture(s, 0, 0, width=Cm(w), height=Cm(h))
ppt.save(r'C:\Users\81901\working\kyopro\kyopro_educational_90/solutions.pptx')    
